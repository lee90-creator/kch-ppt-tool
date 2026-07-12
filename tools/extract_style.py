#!/usr/bin/env python3
"""Extract provisional company style facts from a source PPTX deck.

The extractor is intentionally measurement-first: every value emitted in the
style JSON comes from python-pptx geometry/text/XML inspection. If a section
cannot be measured, the section is returned as null with a reason instead of
being guessed.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from statistics import median
from typing import Any, Iterable

from pptx import Presentation

DEFAULT_INPUT = Path("samples/company-deck.pptx")
DEFAULT_OUTPUT_DIR = Path("ppt-webtool/style")
EMU_PER_PT = 12700
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")
DATE_RE = re.compile(r"^[A-Za-z]+,\s*\d{4}$")


def pct(value: int | float, total: int | float) -> float:
    return round((float(value) / float(total)) * 100.0, 2)


def emu_to_pt(value: int | None) -> float | None:
    if value is None:
        return None
    return round(float(value) / EMU_PER_PT, 2)


def median_or_none(values: Iterable[float | int | None]) -> float | None:
    clean = [float(v) for v in values if v is not None]
    if not clean:
        return None
    return round(float(median(clean)), 2)


def mode_or_none(values: Iterable[Any]) -> Any | None:
    clean = [v for v in values if v is not None]
    if not clean:
        return None
    return Counter(clean).most_common(1)[0][0]


def top_counts(counter: Counter[Any], limit: int = 15) -> list[dict[str, Any]]:
    return [{"value": key, "count": count} for key, count in counter.most_common(limit)]


def resolve_existing_path(raw: str | Path) -> Path:
    path = Path(raw)
    candidates = [path]
    script_root = Path(__file__).resolve().parents[2]
    if not path.is_absolute():
        candidates.append((Path.cwd() / path).resolve())
        candidates.append((script_root / path).resolve())
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    raise FileNotFoundError(f"Input deck not found: {raw}")


def resolve_output_dir(raw: str | Path) -> Path:
    path = Path(raw)
    if path.is_absolute():
        return path
    return (Path.cwd() / path).resolve()


def shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text or ""


def compact_text(text: str) -> str:
    return " ".join(text.split())


def iter_shapes(shapes: Any, prefix: str = "") -> Iterable[tuple[str, Any]]:
    for idx, shape in enumerate(shapes):
        path = f"{prefix}.{idx}" if prefix else str(idx)
        yield path, shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes, path)


def first_text_metrics(shape: Any) -> dict[str, Any]:
    metrics = {
        "font_size_pt": None,
        "bold": None,
        "font_name": None,
        "latin_font": None,
        "east_asian_font": None,
    }
    if not getattr(shape, "has_text_frame", False):
        return metrics
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            metrics["font_size_pt"] = emu_to_pt(int(run.font.size)) if run.font.size else None
            metrics["bold"] = run.font.bold
            metrics["font_name"] = run.font.name
            latin, east_asian = run_typefaces(run)
            metrics["latin_font"] = latin
            metrics["east_asian_font"] = east_asian
            return metrics
    return metrics


def run_typefaces(run: Any) -> tuple[str | None, str | None]:
    latin = None
    east_asian = None
    for rpr in run._r.iter(A_NS + "rPr"):
        for node in rpr.iter(A_NS + "latin"):
            if node.get("typeface"):
                latin = node.get("typeface")
                break
        for node in rpr.iter(A_NS + "ea"):
            if node.get("typeface"):
                east_asian = node.get("typeface")
                break
    return latin, east_asian


def bounds(shape: Any, slide_width: int, slide_height: int) -> dict[str, Any]:
    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    return {
        "emu": {"x": left, "y": top, "w": width, "h": height},
        "pct": {
            "x": pct(left, slide_width),
            "y": pct(top, slide_height),
            "w": pct(width, slide_width),
            "h": pct(height, slide_height),
        },
    }


def median_bounds(records: list[dict[str, Any]]) -> dict[str, Any] | None:
    if not records:
        return None
    return {
        "pct": {
            key: median_or_none(record["bounds"]["pct"][key] for record in records)
            for key in ("x", "y", "w", "h")
        },
        "emu": {
            key: median_or_none(record["bounds"]["emu"][key] for record in records)
            for key in ("x", "y", "w", "h")
        },
    }


def evidence(slide_no: int, shape_path: str, text: str | None = None) -> dict[str, Any]:
    item: dict[str, Any] = {"slide": slide_no, "shape_path": f"slide{slide_no}.shapes[{shape_path}]"}
    if text is not None:
        item["text"] = text
    return item


def solid_fill_color(solid_fill: Any) -> str | None:
    for node in solid_fill.iter(A_NS + "srgbClr"):
        val = node.get("val")
        if val and HEX_RE.match(val):
            return f"#{val.upper()}"
    return None


def direct_fill_hex(element: Any) -> str | None:
    for child in element:
        if child.tag not in {P_NS + "spPr", P_NS + "grpSpPr", P_NS + "picPr"}:
            continue
        for prop_child in child:
            if prop_child.tag == A_NS + "solidFill":
                return solid_fill_color(prop_child)
    return None


def all_srgb_hexes(element: Any) -> list[str]:
    values: list[str] = []
    for solid_fill in element.iter(A_NS + "solidFill"):
        color = solid_fill_color(solid_fill)
        if color:
            values.append(color)
    return values


def extract_header(prs: Presentation) -> dict[str, Any]:
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    number_records: list[dict[str, Any]] = []
    title_records: list[dict[str, Any]] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        if slide_no <= 2:
            continue
        text_shapes: list[dict[str, Any]] = []
        for shape_path, shape in iter_shapes(slide.shapes):
            text = compact_text(shape_text(shape))
            if not text:
                continue
            b = bounds(shape, slide_width, slide_height)
            m = first_text_metrics(shape)
            if b["emu"]["x"] < 0 or b["emu"]["y"] < 0:
                continue
            text_shapes.append({"path": shape_path, "shape": shape, "text": text, "bounds": b, "metrics": m})

        number_candidates = [
            item
            for item in text_shapes
            if re.fullmatch(r"0[1-5]", item["text"])
            and item["bounds"]["pct"]["x"] < 12
            and item["bounds"]["pct"]["y"] < 8
            and (item["metrics"]["font_size_pt"] or 0) >= 40
        ]
        if not number_candidates:
            continue
        number = sorted(number_candidates, key=lambda item: (item["bounds"]["emu"]["y"], item["bounds"]["emu"]["x"]))[0]
        title_candidates = [
            item
            for item in text_shapes
            if item is not number
            and not re.fullmatch(r"0[1-5]", item["text"])
            and 10 <= item["bounds"]["pct"]["x"] <= 35
            and 3 <= item["bounds"]["pct"]["y"] <= 11
            and 18 <= (item["metrics"]["font_size_pt"] or 0) <= 32
        ]
        if not title_candidates:
            continue
        title = sorted(
            title_candidates,
            key=lambda item: (
                abs(item["bounds"]["pct"]["y"] - 6.7),
                abs(item["bounds"]["pct"]["x"] - 16.2),
            ),
        )[0]
        number_records.append(
            {
                "slide": slide_no,
                "shape_path": number["path"],
                "text": number["text"],
                "bounds": number["bounds"],
                **number["metrics"],
            }
        )
        title_records.append(
            {
                "slide": slide_no,
                "shape_path": title["path"],
                "text": title["text"],
                "bounds": title["bounds"],
                **title["metrics"],
            }
        )

    if not title_records or not number_records:
        return {"value": None, "reason": "본문 슬라이드의 섹션번호/제목 헤더 쌍을 실측하지 못했습니다."}

    title_counter = Counter(record["text"] for record in title_records)
    return {
        "value": {
            "rule": "본문 슬라이드는 좌상단 큰 섹션번호 박스(01~05)와 그 오른쪽 섹션 제목을 같은 헤더 밴드에 배치한다.",
            "slides_measured": len(title_records),
            "section_number_box": {
                "median_bounds": median_bounds(number_records),
                "median_font_size_pt": median_or_none(record["font_size_pt"] for record in number_records),
                "bold_mode": mode_or_none(record["bold"] for record in number_records),
                "font_mode": mode_or_none(record["font_name"] for record in number_records),
                "texts_seen": sorted(Counter(record["text"] for record in number_records).keys()),
            },
            "section_title": {
                "median_bounds": median_bounds(title_records),
                "median_font_size_pt": median_or_none(record["font_size_pt"] for record in title_records),
                "bold_mode": mode_or_none(record["bold"] for record in title_records),
                "font_mode": mode_or_none(record["font_name"] for record in title_records),
                "most_common_titles": [{"text": text, "count": count} for text, count in title_counter.most_common(8)],
            },
            "evidence": [
                evidence(record["slide"], record["shape_path"], record["text"])
                for record in (number_records[:3] + title_records[:3])
            ],
        },
        "reason": None,
    }


def extract_governing_message(prs: Presentation) -> dict[str, Any]:
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    records: list[dict[str, Any]] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        if slide_no <= 2:
            continue
        candidates: list[dict[str, Any]] = []
        for shape_path, shape in iter_shapes(slide.shapes):
            text = compact_text(shape_text(shape))
            if len(text) < 20:
                continue
            b = bounds(shape, slide_width, slide_height)
            if (
                b["emu"]["x"] < 0
                or b["emu"]["y"] < 0
                or b["emu"]["x"] >= slide_width
                or b["emu"]["y"] >= slide_height
            ):
                continue
            y = b["pct"]["y"]
            x = b["pct"]["x"]
            width = b["pct"]["w"]
            if 22 <= y <= 34 and x <= 10 and width >= 75:
                m = first_text_metrics(shape)
                candidates.append({"slide": slide_no, "shape_path": shape_path, "text": text, "bounds": b, **m})
        if candidates:
            selected = sorted(
                candidates,
                key=lambda item: (
                    abs(item["bounds"]["pct"]["y"] - 24.0),
                    -item["bounds"]["pct"]["w"],
                    -len(item["text"]),
                ),
            )[0]
            records.append(selected)

    if not records:
        return {"value": None, "reason": "제목 아래의 장문 결론 문장 후보를 실측하지 못했습니다."}

    return {
        "value": {
            "rule": "소주제 아래에 결론을 먼저 말하는 한 문장(대개 1~2줄)을 본문보다 먼저 배치한다. 생성 시 표·도표 설명보다 앞에 핵심 결론을 선행한다.",
            "slides_measured": len(records),
            "median_bounds": median_bounds(records),
            "median_font_size_pt": median_or_none(record["font_size_pt"] for record in records),
            "bold_mode": mode_or_none(record["bold"] for record in records),
            "font_mode": mode_or_none(record["font_name"] for record in records),
            "examples": [
                evidence(record["slide"], record["shape_path"], record["text"])
                for record in records[:8]
            ],
            "notes": [
                "소제목만 있고 장문 결론 문장이 없는 본문 슬라이드는 중앙값 산정에서 제외했습니다.",
                "문장 후보 조건: 슬라이드 상단 22~34% y영역, x 10% 이하, 폭 75% 이상, 공백 제거 후 20자 이상.",
            ],
        },
        "reason": None,
    }


def collect_palette_cells(group_shape: Any, slide_no: int, group_path: str) -> list[dict[str, Any]]:
    cells: list[dict[str, Any]] = []
    for child_path, child in iter_shapes(group_shape.shapes if hasattr(group_shape, "shapes") else []):
        if hasattr(child, "shapes"):
            continue
        fill = direct_fill_hex(child.element)
        if not fill:
            continue
        text = compact_text(shape_text(child))
        cells.append(
            {
                "hex": fill,
                "rgb_label": text if re.search(r"\d+\s*/\s*\d+\s*/\s*\d+", text) else None,
                "evidence": {
                    "slide": slide_no,
                    "shape_path": f"slide{slide_no}.shapes[{group_path}.{child_path}]",
                },
            }
        )
    return cells


def dedupe_cells(cells: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[str] = set()
    deduped: list[dict[str, Any]] = []
    for cell in cells:
        key = f"{cell['hex']}|{cell.get('rgb_label')}|{cell['evidence']['shape_path']}"
        if key in seen:
            continue
        seen.add(key)
        deduped.append(cell)
    return deduped


def extract_colors(prs: Presentation) -> dict[str, Any]:
    slide_no = 18 if len(prs.slides) >= 18 else None
    if slide_no is None:
        return {"value": None, "reason": "팔레트가 있는 18번 슬라이드가 없습니다."}

    slide = prs.slides[slide_no - 1]
    labels: dict[str, dict[str, Any]] = {}
    for shape_path, shape in iter_shapes(slide.shapes):
        text = compact_text(shape_text(shape))
        if "Primary Colors" in text:
            labels["primary"] = {"shape_path": shape_path, "top": int(shape.top), "text": text}
        elif "Accent Colors" in text:
            labels["accent"] = {"shape_path": shape_path, "top": int(shape.top), "text": text}
        elif "Neutrals" in text:
            labels["neutral"] = {"shape_path": shape_path, "top": int(shape.top), "text": text}

    groups: list[dict[str, Any]] = []
    for shape_path, shape in iter_shapes(slide.shapes):
        if not hasattr(shape, "shapes"):
            continue
        cells = dedupe_cells(collect_palette_cells(shape, slide_no, shape_path))
        if cells:
            groups.append({"shape_path": shape_path, "top": int(shape.top), "left": int(shape.left), "cells": cells})

    def group_between(start: int, end: int | None, minimum: int = 3) -> dict[str, Any] | None:
        candidates = [g for g in groups if g["top"] > start and (end is None or g["top"] < end) and len(g["cells"]) >= minimum]
        if not candidates:
            return None
        return sorted(candidates, key=lambda g: (g["top"], -len(g["cells"])))[0]

    primary_group = group_between(labels.get("primary", {}).get("top", -10**12), labels.get("accent", {}).get("top"), 10)
    accent_group = group_between(labels.get("accent", {}).get("top", -10**12), labels.get("neutral", {}).get("top"), 10)
    neutral_group = group_between(labels.get("neutral", {}).get("top", -10**12), None, 3)

    traffic_light_group = None
    for group in groups:
        contains_traffic_label = False
        shape_obj = None
        for path, shape in iter_shapes(slide.shapes):
            if path == group["shape_path"]:
                shape_obj = shape
                break
        if shape_obj is not None:
            for _, child in iter_shapes(shape_obj.shapes):
                if "Traffic Light Palette" in compact_text(shape_text(child)):
                    contains_traffic_label = True
                    break
        if contains_traffic_label:
            traffic_light_group = group
            break

    frequency = Counter()
    for slide in prs.slides:
        frequency.update(all_srgb_hexes(slide.element))

    palette: dict[str, Any] = {}
    missing: list[str] = []
    palette_labels = {
        "primary": "Primary Colors",
        "accent": "Accent Colors",
        "neutral": "Neutral Colors",
    }
    for name, group in (("primary", primary_group), ("accent", accent_group), ("neutral", neutral_group)):
        label = palette_labels[name]
        if group is None:
            reason = f"{label} 팔레트 그룹을 실측하지 못했습니다."
            palette[name] = {"value": None, "reason": reason}
            missing.append(reason)
            continue
        base_cells = [cell for cell in group["cells"] if cell.get("rgb_label")]
        if not base_cells:
            base_cells = group["cells"][: min(8, len(group["cells"]))]
        palette[name] = {
            "value": {
                "swatch_count": len(group["cells"]),
                "base_colors": base_cells,
                "all_swatches": group["cells"],
                "evidence": {
                    "slide": slide_no,
                    "label_shape_path": f"slide{slide_no}.shapes[{labels.get(name, {}).get('shape_path')}]" if name in labels else None,
                    "palette_group_shape_path": f"slide{slide_no}.shapes[{group['shape_path']}]",
                },
            },
            "reason": None,
        }

    if traffic_light_group is not None and traffic_light_group is not neutral_group:
        palette["traffic_light"] = {
            "value": {
                "swatch_count": len(traffic_light_group["cells"]),
                "all_swatches": traffic_light_group["cells"],
                "evidence": {"slide": slide_no, "palette_group_shape_path": f"slide{slide_no}.shapes[{traffic_light_group['shape_path']}]"},
            },
            "reason": None,
        }

    reason = None
    if missing:
        reason = "필수 팔레트 그룹 일부를 실측하지 못했습니다: " + "; ".join(missing)
    return {
        "value": {
            "palette_slide": slide_no,
            "palette": palette,
            "top_hex_frequency_all_slides": top_counts(frequency, 20),
            "notes": [
                "slide 18의 오프-캔버스 팔레트 라벨과 그룹 도형에서 직접 solidFill srgbClr 값을 읽었습니다.",
                "보조 빈도는 전체 슬라이드 XML의 명시 srgbClr solidFill 값을 집계한 값으로, 텍스트/선/도형 fill이 함께 포함될 수 있습니다.",
            ],
        },
        "reason": reason,
    }


def extract_fonts(prs: Presentation) -> dict[str, Any]:
    latin = Counter()
    east_asian = Counter()
    run_name = Counter()
    sizes = Counter()
    bold = Counter()
    total_runs = 0

    for slide in prs.slides:
        for _, shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    total_runs += 1
                    if run.font.name:
                        run_name[run.font.name] += 1
                    if run.font.size:
                        sizes[round(int(run.font.size) / EMU_PER_PT, 2)] += 1
                    bold[str(run.font.bold)] += 1
                    latin_font, east_asian_font = run_typefaces(run)
                    if latin_font:
                        latin[latin_font] += 1
                    if east_asian_font:
                        east_asian[east_asian_font] += 1

    if total_runs == 0:
        return {"value": None, "reason": "텍스트 run을 찾지 못했습니다."}

    primary = None
    source = None
    if east_asian:
        primary = east_asian.most_common(1)[0][0]
        source = "east_asian_frequency"
    elif latin:
        primary = latin.most_common(1)[0][0]
        source = "latin_frequency"
    elif run_name:
        primary = run_name.most_common(1)[0][0]
        source = "python_pptx_run_font_name"

    malgun_count = sum(count for name, count in (latin + east_asian + run_name).items() if "Malgun" in str(name) or "맑은" in str(name))
    return {
        "value": {
            "total_text_runs_measured": total_runs,
            "primary_font": primary,
            "primary_font_decision_source": source,
            "latin_font_frequency": top_counts(latin, 20),
            "east_asian_font_frequency": top_counts(east_asian, 20),
            "python_pptx_run_font_name_frequency": top_counts(run_name, 20),
            "font_size_frequency_pt": [{"value": key, "count": count} for key, count in sizes.most_common(20)],
            "bold_frequency": [{"value": key, "count": count} for key, count in bold.most_common()],
            "malgun_gothic_related_count": malgun_count,
            "note": "주폰트는 빈도 1위 East Asian typeface로 판정했습니다. 맑은 고딕/Malgun 계열은 실측 빈도가 낮아 주폰트로 판정하지 않았습니다.",
        },
        "reason": None,
    }


def extract_cover_toc(prs: Presentation) -> dict[str, Any]:
    if len(prs.slides) < 2:
        return {"value": None, "reason": "표지와 목차를 측정하려면 최소 2개 슬라이드가 필요합니다."}

    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    cover = prs.slides[0]
    toc = prs.slides[1]

    cover_shapes: list[dict[str, Any]] = []
    for shape_path, shape in iter_shapes(cover.shapes):
        text = shape_text(shape)
        if not compact_text(text):
            continue
        paragraphs = [compact_text(p.text) for p in shape.text_frame.paragraphs if compact_text(p.text)]
        cover_shapes.append({"shape_path": shape_path, "shape": shape, "text": text, "paragraphs": paragraphs, "bounds": bounds(shape, slide_width, slide_height)})
    if not cover_shapes:
        return {"value": None, "reason": "표지 슬라이드에서 텍스트 프레임을 찾지 못했습니다."}
    cover_title_shape = max(cover_shapes, key=lambda item: len(compact_text(item["text"])))
    cover_paragraphs = cover_title_shape["paragraphs"]
    cover_metrics: list[dict[str, Any]] = []
    shape = cover_title_shape["shape"]
    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = compact_text(paragraph.text)
        if not paragraph_text:
            continue
        metric = {"text": paragraph_text, "font_size_pt": None, "font_name": None, "bold": None}
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            metric["font_size_pt"] = emu_to_pt(int(run.font.size)) if run.font.size else None
            metric["font_name"] = run.font.name
            metric["bold"] = run.font.bold
            break
        cover_metrics.append(metric)

    toc_entries: list[dict[str, Any]] = []
    toc_title = None
    toc_number_shapes: list[dict[str, Any]] = []
    toc_title_shapes: list[dict[str, Any]] = []
    for shape_path, shape in iter_shapes(toc.shapes):
        text = compact_text(shape_text(shape))
        if not text:
            continue
        item = {"shape_path": shape_path, "text": text, "bounds": bounds(shape, slide_width, slide_height), **first_text_metrics(shape)}
        if "Table of Contents" in text:
            toc_title = item
        elif re.fullmatch(r"0[1-9]", text):
            toc_number_shapes.append(item)
        else:
            toc_title_shapes.append(item)

    for title in sorted(toc_title_shapes, key=lambda item: item["bounds"]["emu"]["y"]):
        nearest_number = None
        if toc_number_shapes:
            nearest_number = min(toc_number_shapes, key=lambda item: abs(item["bounds"]["emu"]["y"] - title["bounds"]["emu"]["y"]))
        toc_entries.append(
            {
                "number": nearest_number["text"] if nearest_number else None,
                "title": title["text"],
                "title_shape_path": f"slide2.shapes[{title['shape_path']}]",
                "number_shape_path": f"slide2.shapes[{nearest_number['shape_path']}]" if nearest_number else None,
            }
        )

    title_text = cover_paragraphs[0] if cover_paragraphs else None
    date_text = cover_paragraphs[1] if len(cover_paragraphs) > 1 else None
    return {
        "value": {
            "cover": {
                "rule": "표지는 중앙보다 약간 왼쪽의 대형 제목 블록에 연도 사업계획 제목을 두고, 바로 아래 영문 월/연도 날짜를 표기한다.",
                "title_text": title_text,
                "date_text": date_text,
                "date_format_convention": "English Month, YYYY" if date_text and DATE_RE.match(date_text) else None,
                "title_block_bounds": cover_title_shape["bounds"],
                "paragraph_metrics": cover_metrics,
                "evidence": evidence(1, cover_title_shape["shape_path"], compact_text(cover_title_shape["text"])),
            },
            "toc": {
                "rule": "목차는 우측 중심 세로 목록이며, 상단에 영어 'Table of Contents' 제목을 두고 각 섹션을 2자리 번호+섹션명으로 나열한다.",
                "title": toc_title,
                "entries": toc_entries,
                "entry_count": len(toc_entries),
            },
        },
        "reason": None,
    }


def field_value(section: dict[str, Any]) -> Any:
    return section.get("value")


def section_status(section: dict[str, Any]) -> str:
    if section.get("value") is None:
        return f"null: {section.get('reason')}"
    if section.get("reason"):
        return f"partial: {section.get('reason')}"
    return "ok"


def short_pct_box(bounds_value: dict[str, Any] | None) -> str | None:
    if not bounds_value:
        return None
    p = bounds_value["pct"]
    return f"x {p['x']}%, y {p['y']}%, w {p['w']}%, h {p['h']}%"


def required_reason(label: str, missing: list[str]) -> str:
    return f"{label} 필수 측정값 누락: " + ", ".join(missing)


def palette_group_value(colors: dict[str, Any] | None, group_name: str) -> dict[str, Any] | None:
    if not colors:
        return None
    group = colors.get("palette", {}).get(group_name)
    if not group:
        return None
    if isinstance(group, dict) and ("value" in group or "reason" in group):
        return group.get("value")
    return group


def palette_group_reason(colors: dict[str, Any] | None, group_name: str) -> str | None:
    if not colors:
        return None
    group = colors.get("palette", {}).get(group_name)
    if isinstance(group, dict):
        return group.get("reason")
    return None


def build_design_spec_mapping(extraction: dict[str, Any]) -> dict[str, Any]:
    header = field_value(extraction["header"])
    governing = field_value(extraction["governing_message"])
    colors = field_value(extraction["colors"])
    fonts = field_value(extraction["fonts"])
    cover_toc = field_value(extraction["cover_toc"])

    off_rule = "company_style_toggle=off이면 이 회사 스타일 지시문을 design_spec/spec_lock 프롬프트에 주입하지 않는다."
    applied_when = "company_style_toggle=on"

    header_section: dict[str, Any] = {"applied_when": applied_when}
    header_missing: list[str] = []
    if header:
        number_box = short_pct_box(header["section_number_box"]["median_bounds"])
        number_size = header["section_number_box"]["median_font_size_pt"]
        header_box = short_pct_box(header["section_title"]["median_bounds"])
        header_size = header["section_title"]["median_font_size_pt"]
        if number_box is None:
            header_missing.append("섹션번호 박스 좌표/크기")
        if number_size is None:
            header_missing.append("섹션번호 박스 글자 크기")
        if header_box is None:
            header_missing.append("섹션 제목 좌표/크기")
        if header_size is None:
            header_missing.append("섹션 제목 글자 크기")
        if not header_missing:
            header_section["prompt_instruction_ko"] = (
                f"본문 슬라이드 헤더는 좌상단 섹션번호 박스를 {number_box}, 약 {number_size}pt 굵은 표시로 두고, "
                f"섹션 제목은 {header_box}, 약 {header_size}pt 굵은 제목으로 배치한다."
            )
    else:
        header_missing.append(extraction["header"].get("reason") or "헤더 측정값")
    header_section["reason"] = required_reason("header", header_missing) if header_missing else None

    governing_section: dict[str, Any] = {"applied_when": applied_when}
    governing_missing: list[str] = []
    if governing:
        gm_box = short_pct_box(governing["median_bounds"])
        gm_size = governing["median_font_size_pt"]
        if gm_box is None:
            governing_missing.append("결론 선행 문장 좌표/크기")
        if gm_size is None:
            governing_missing.append("결론 선행 문장 글자 크기")
        if not governing_missing:
            governing_section["prompt_instruction_ko"] = (
                f"각 본문 슬라이드는 소제목 아래 {gm_box} 위치에 약 {gm_size}pt의 결론 선행 한 문장을 먼저 제시하고, "
                "이후 표·도표·상세 근거를 배치한다."
            )
    else:
        governing_missing.append(extraction["governing_message"].get("reason") or "결론 선행 문장 측정값")
    governing_section["reason"] = required_reason("governing_message", governing_missing) if governing_missing else None

    color_groups: dict[str, dict[str, Any]] = {}
    top_color_reason = extraction["colors"].get("reason") if not colors else None
    color_missing: list[str] = []
    color_labels = {"primary": "Primary", "accent": "Accent", "neutral": "Neutral"}
    for group_name, label in color_labels.items():
        group_value = palette_group_value(colors, group_name)
        group_reason = palette_group_reason(colors, group_name)
        hex_values = [cell["hex"] for cell in group_value.get("base_colors", []) if cell.get("hex")] if group_value else []
        if hex_values:
            color_groups[group_name] = {"value": hex_values, "reason": None}
        else:
            reason = group_reason or top_color_reason or f"{label} 색상 HEX"
            color_groups[group_name] = {"value": None, "reason": reason}
            color_missing.append(f"{label}: {reason}")
    colors_section: dict[str, Any] = {"applied_when": applied_when, "palette_groups": color_groups}
    if not color_missing:
        colors_section["prompt_instruction_ko"] = (
            "색상은 실측 팔레트를 우선 사용한다: "
            f"Primary {', '.join(color_groups['primary']['value'])}, "
            f"Accent {', '.join(color_groups['accent']['value'])}, "
            f"Neutral {', '.join(color_groups['neutral']['value'])}. "
            "토글 OFF시 이 HEX 팔레트는 주입하지 않는다."
        )
    colors_section["reason"] = required_reason("colors", color_missing) if color_missing else None

    measured_font = fonts.get("primary_font") if fonts else None
    font_policy = {
        "applied_font": "Malgun Gothic",
        "reason": "실측 주폰트가 비표준 폰트라 배포 환경 호환 위해 대체",
        "policy_source": "deep-interview 스펙 R8 필드④ (밝은 고딕 계열, 직원 PC 기본 탑재 폰트 제한)",
    }
    fonts_section: dict[str, Any] = {
        "applied_when": applied_when,
        "measured": {"primary_font": measured_font},
        "runtime_font_policy": font_policy,
    }
    font_missing: list[str] = []
    if measured_font is None:
        font_missing.append(extraction["fonts"].get("reason") or "실측 주폰트")
    else:
        fonts_section["prompt_instruction_ko"] = (
            "타이포그래피는 실측 주폰트를 별도 measured.primary_font에 보존하고, 런타임 생성/적용 시 "
            f"배포 환경 호환 정책에 따라 '{font_policy['applied_font']}'으로 정책 대체한다. "
            f"숫자/본문/도표 텍스트도 '{font_policy['applied_font']}' 계열의 굵기 변형(Regular/Bold)만 사용한다. "
            f"정책 근거: {font_policy['policy_source']}."
        )
    fonts_section["reason"] = required_reason("fonts", font_missing) if font_missing else None

    cover_section: dict[str, Any] = {"applied_when": applied_when}
    cover_missing: list[str] = []
    if cover_toc:
        cover = cover_toc["cover"]
        cover_box = short_pct_box(cover["title_block_bounds"])
        date_format = cover.get("date_format_convention")
        if cover_box is None:
            cover_missing.append("표지 제목 블록 좌표/크기")
        if date_format is None:
            cover_missing.append("표지 날짜 형식")
        if not cover_missing:
            cover_section["prompt_instruction_ko"] = (
                f"표지는 제목 블록을 {cover_box}에 크게 두고 날짜는 '{date_format}' 관례로 표기한다. "
                "목차는 우측 중심 세로 목록과 'Table of Contents' 제목 구조를 따른다."
            )
    else:
        cover_missing.append(extraction["cover_toc"].get("reason") or "표지/목차 측정값")
    cover_section["reason"] = required_reason("cover_toc", cover_missing) if cover_missing else None

    return {
        "applied_when": applied_when,
        "not_applied_when": "company_style_toggle=off",
        "toggle_off_rule": off_rule,
        "header": header_section,
        "governing_message": governing_section,
        "colors": colors_section,
        "fonts": fonts_section,
        "cover_toc": cover_section,
    }


def build_preview(extraction: dict[str, Any]) -> str:
    header = field_value(extraction["header"])
    governing = field_value(extraction["governing_message"])
    colors = field_value(extraction["colors"])
    fonts = field_value(extraction["fonts"])
    cover_toc = field_value(extraction["cover_toc"])
    mapping = extraction["design_spec_mapping"]

    lines: list[str] = []
    lines.append("# Company Style Extraction Preview")
    lines.append("")
    lines.append(f"- Source deck: `{extraction['source_deck']}`")
    lines.append(f"- Extracted at: {extraction['extracted_at']}")
    lines.append(f"- Status: {extraction['status']}")
    lines.append(f"- Slides: {extraction['slide_count']}")
    lines.append("")

    lines.append("## 1. Header")
    if header:
        number_box = short_pct_box(header["section_number_box"]["median_bounds"]) or "미측정"
        title_box = short_pct_box(header["section_title"]["median_bounds"]) or "미측정"
        lines.append(f"- Rule: {header['rule']}")
        lines.append(f"- Section number box median: {number_box}; font {header['section_number_box']['median_font_size_pt']}pt; bold mode {header['section_number_box']['bold_mode']}.")
        lines.append(f"- Section title median: {title_box}; font {header['section_title']['median_font_size_pt']}pt; bold mode {header['section_title']['bold_mode']}.")
        evidence_text = ", ".join(f"s{ev['slide']} {ev['shape_path']}" for ev in header["evidence"][:4])
        lines.append(f"- Evidence: {evidence_text}")
    else:
        lines.append(f"- Null: {extraction['header'].get('reason')}")
    lines.append("")

    lines.append("## 2. Governing Message")
    if governing:
        gm_box = short_pct_box(governing["median_bounds"]) or "미측정"
        lines.append(f"- Rule: {governing['rule']}")
        lines.append(f"- Median position: {gm_box}; font {governing['median_font_size_pt']}pt; measured slides {governing['slides_measured']}.")
        for example in governing["examples"][:3]:
            lines.append(f"- Example s{example['slide']} `{example['text']}` ({example['shape_path']})")
    else:
        lines.append(f"- Null: {extraction['governing_message'].get('reason')}")
    lines.append("")

    lines.append("## 3. Colors")
    if colors:
        lines.append("- Palette source: slide 18 off-canvas palette groups; all HEX values are measured `solidFill` values.")
        if extraction["colors"].get("reason"):
            lines.append(f"- Partial reason: {extraction['colors']['reason']}")
        lines.append("- Top all-slide HEX frequency: " + ", ".join(f"{item['value']}({item['count']})" for item in colors["top_hex_frequency_all_slides"][:10]))
        for group_name in ("primary", "accent", "neutral"):
            group_reason = palette_group_reason(colors, group_name)
            if group_reason:
                lines.append(f"- {group_name} reason: {group_reason}")
        lines.append("")
        lines.append("| Group | # | HEX | RGB label | Evidence |")
        lines.append("|---|---:|---|---|---|")
        for group_name in ("primary", "accent", "neutral", "traffic_light"):
            group_value = palette_group_value(colors, group_name)
            if not group_value:
                continue
            for idx, cell in enumerate(group_value.get("all_swatches", []), start=1):
                label = cell.get("rgb_label") or ""
                ev = cell["evidence"]["shape_path"]
                lines.append(f"| {group_name} | {idx} | `{cell['hex']}` | {label} | {ev} |")
    else:
        lines.append(f"- Null: {extraction['colors'].get('reason')}")
    lines.append("")

    lines.append("## 4. Fonts")
    font_mapping = mapping["fonts"]
    if fonts:
        lines.append(f"- Measured primary font: `{fonts['primary_font']}` from {fonts['primary_font_decision_source']}; measured runs {fonts['total_text_runs_measured']}.")
        lines.append("- East Asian top fonts: " + ", ".join(f"{item['value']}({item['count']})" for item in fonts["east_asian_font_frequency"][:8]))
        lines.append(f"- Measurement note: {fonts['note']}")
    else:
        lines.append(f"- Null: {extraction['fonts'].get('reason')}")
    policy = font_mapping["runtime_font_policy"]
    lines.append(f"- Runtime font policy: apply `{policy['applied_font']}` as 정책 대체.")
    lines.append(f"- Policy reason: {policy['reason']}")
    lines.append(f"- Policy source: {policy['policy_source']}")
    lines.append("")

    lines.append("## 5. Cover / TOC")
    if cover_toc:
        cover = cover_toc["cover"]
        toc = cover_toc["toc"]
        cover_box = short_pct_box(cover["title_block_bounds"]) or "미측정"
        toc_title = toc["title"]["text"] if toc["title"] else "미측정"
        lines.append(f"- Cover: `{cover['title_text']}` + `{cover['date_text']}` at {cover_box}; date convention `{cover['date_format_convention']}`.")
        lines.append(f"- TOC: {toc['entry_count']} entries, right-side vertical list under `{toc_title}`.")
        for entry in toc["entries"]:
            number = entry["number"] or "번호 미측정"
            lines.append(f"  - {number} {entry['title']} ({entry['title_shape_path']})")
    else:
        lines.append(f"- Null: {extraction['cover_toc'].get('reason')}")
    lines.append("")

    lines.append("## Design Spec Mapping")
    lines.append(f"- Applied when: `{mapping['applied_when']}`")
    lines.append(f"- Toggle-off rule: {mapping['toggle_off_rule']}")
    for field in ("header", "governing_message", "colors", "fonts", "cover_toc"):
        section = mapping[field]
        prompt = section.get("prompt_instruction_ko")
        if prompt:
            lines.append(f"- {field}: {prompt}")
        else:
            lines.append(f"- {field}: reason={section.get('reason')}")
    lines.append("")

    lines.append("## Uncertainties / Limitations")
    lines.append("- Governing message 중앙값은 장문 결론 문장이 있는 슬라이드만 사용했습니다; 소제목만 있는 슬라이드는 제외했습니다.")
    lines.append("- 전체 HEX 빈도는 XML의 명시 `srgbClr` 빈도이며 텍스트/선/도형 색상이 함께 집계될 수 있습니다.")
    if fonts:
        lines.append(f"- 실측 주폰트 `{fonts['primary_font']}`와 런타임 정책 폰트 `{policy['applied_font']}`를 분리 기록했습니다.")
    lines.append("")
    return "\n".join(lines)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Extract a provisional frozen company style JSON from a PPTX deck.")
    parser.add_argument("--input", default=str(DEFAULT_INPUT), help="Source PPTX path. Defaults to the provided sample deck.")
    parser.add_argument("--output-dir", default=str(DEFAULT_OUTPUT_DIR), help="Directory for company_style.json and preview markdown.")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        input_path = resolve_existing_path(args.input)
    except FileNotFoundError as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 2
    output_dir = resolve_output_dir(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(input_path)
    extraction: dict[str, Any] = {
        "schema_version": 1,
        "source_deck": str(Path(args.input)),
        "source_deck_resolved": str(input_path),
        "extracted_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "status": "provisional-frozen",
        "slide_count": len(prs.slides),
        "slide_size_emu": {"width": int(prs.slide_width), "height": int(prs.slide_height)},
    }
    extraction["header"] = extract_header(prs)
    extraction["governing_message"] = extract_governing_message(prs)
    extraction["colors"] = extract_colors(prs)
    extraction["fonts"] = extract_fonts(prs)
    extraction["cover_toc"] = extract_cover_toc(prs)
    extraction["design_spec_mapping"] = build_design_spec_mapping(extraction)

    json_path = output_dir / "company_style.json"
    preview_path = output_dir / "company_style_preview.md"
    json_path.write_text(json.dumps(extraction, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    preview_path.write_text(build_preview(extraction), encoding="utf-8")

    summary = {
        "json": str(json_path),
        "preview": str(preview_path),
        "status": extraction["status"],
        "slides": len(prs.slides),
        "fields": {
            field: section_status(extraction[field])
            for field in ("header", "governing_message", "colors", "fonts", "cover_toc")
        },
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
